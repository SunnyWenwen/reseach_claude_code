from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x22, 0x2E)
MUTED  = RGBColor(0x55, 0x65, 0x7A)
BLUE   = RGBColor(0x1D, 0x6F, 0xC8)
GRAY   = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT  = RGBColor(0xF8, 0xFA, 0xFC)
GREEN  = RGBColor(0x05, 0x7A, 0x50)
ORANGE = RGBColor(0xB4, 0x5A, 0x00)
RED    = RGBColor(0xAA, 0x18, 0x18)
PURPLE = RGBColor(0x5B, 0x21, 0xB6)
TEAL   = RGBColor(0x0E, 0x7A, 0x8E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────

def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def hline(slide, y, x=0.5, w=12.33, color=GRAY, h=0.045):
    bar = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def rect(slide, x, y, w, h, color=GRAY):
    r = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def t(slide, text, x, y, w, h=0.55,
      size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def slide_header(slide, title, subtitle=None):
    bg(slide)
    hline(slide, 0.0, x=0.0, w=13.33, color=BLUE, h=0.10)
    t(slide, title, 0.55, 0.18, 12.2, 0.7, size=26, bold=True, color=INK)
    hline(slide, 1.0, color=GRAY)
    if subtitle:
        t(slide, subtitle, 0.55, 1.08, 12, 0.42, size=14, color=MUTED, italic=True)

def bullets(slide, items, x, y, w, size=17, gap=0.48, color=INK):
    cy = y
    for item in items:
        is_sub = item.startswith(">>")
        text = ("    • " + item[2:]) if is_sub else ("• " + item)
        c = MUTED if is_sub else color
        sz = size - 2 if is_sub else size
        t(slide, text, x, cy, w, gap, size=sz, color=c)
        cy += gap
    return cy

def label(slide, text, x, y, color=BLUE, size=13):
    t(slide, text.upper(), x, y, 6, 0.33, size=size, bold=True, color=color)

def reason_slide(num, col, title, subtitle, pts):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    t(s, f"Reason {num}", 0.65, 0.38, 3.5, 0.38, size=15, bold=True, color=col)
    t(s, title,    0.65, 0.70, 11.7, 0.88, size=36, bold=True, color=INK)
    t(s, subtitle, 0.65, 1.55, 11.5, 0.48, size=18, color=MUTED, italic=True)
    hline(s, 2.1, x=0.65, w=12.03)
    cy = 2.3
    for pt in pts:
        t(s, "• " + pt, 0.65, cy, 12.1, 0.88, size=20, color=INK)
        cy += 0.92
    return s


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 – Cover
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.33, 4.2, color=BLUE)
t(s, "Claude Code", 0.8, 0.7, 11.5, 1.4, size=60, bold=True, color=WHITE)
t(s, "Why It's Strong & How It Works",
  0.8, 2.15, 11, 0.75, size=26, color=RGBColor(0xBB, 0xD6, 0xF8))
hline(s, 4.22, x=0, w=13.33, color=BLUE, h=0.06)
t(s, "Hsiang-Wen", 0.8, 4.65, 5, 0.5, size=20, color=INK)
t(s, "2026",       0.8, 5.22, 3, 0.4, size=16, color=MUTED)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 – Agenda
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Agenda")

sections = [
    ("Part 1", "Why Claude Code 特別強", BLUE, [
        ("01", "4 個核心原因：Agentic Loop、Context Management、Strong Tools、Strong Model"),
        ("02", "案例帶入：折扣碼重複使用的 Bug 修復"),
    ]),
    ("Part 2", "How It Works  機制深入", GREEN, [
        ("03", "Agent Architecture — ReAct Loop"),
        ("04", "工具總覽 — 依功能分類"),
        ("05", "工具深入：Read / Grep / Glob、Edit / Write、Bash、ToolSearch"),
        ("06", "工具深入：Agent、Skill"),
        ("07", "系統機制：Hooks、Permission"),
    ]),
    ("Part 3", "Take Away", ORANGE, [
        ("08", "Takeaways"),
    ]),
]
cy = 1.22
for part_label, part_title, col, items in sections:
    hline(s, cy, color=GRAY)
    rect(s, 0.5, cy, 12.33, 0.60, color=col)
    t(s, part_label, 0.65, cy+0.12, 1.2, 0.38, size=13, bold=True, color=WHITE)
    t(s, part_title, 1.88, cy+0.12, 10.8, 0.38, size=16, bold=True, color=WHITE)
    cy += 0.62
    for num, title in items:
        hline(s, cy, color=GRAY)
        t(s, num,   0.75, cy+0.10, 0.75, 0.40, size=16, bold=True, color=col)
        t(s, title, 1.58, cy+0.10, 11.1, 0.40, size=16, color=INK)
        cy += 0.58


# ══════════════════════════════════════════════════════════════════
# SLIDES 3–7 – Why: Overall + 4 Reasons
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
t(s, "Claude Code 為何特別強？",
  0.65, 0.48, 11.5, 0.95, size=40, bold=True)
t(s, "不是單點能力，而是完整的 Agentic Coding System",
  0.65, 1.38, 11.5, 0.44, size=18, color=MUTED, italic=True)
hline(s, 1.92, x=0.65, w=12.03)
boxes = [
    (BLUE,   "01  Agentic Loop",
     "工作方式是循環式：理解 → 行動 → 驗證，能持續根據結果修正"),
    (GREEN,  "02  Context Management",
     "長任務過程中主動維持 context 品質，避免雜訊讓模型失焦"),
    (ORANGE, "03  Strong Tools",
     "真正能搜尋、讀檔、改檔、執行驗證，不只是給建議"),
    (PURPLE, "04  Strong Model",
     "理解需求、規劃步驟、選對工具、解讀結果，並依任務分配最適模型"),
]
grid = [(0.65, 2.05), (7.0, 2.05), (0.65, 4.45), (7.0, 4.45)]
BW, BH = 5.9, 2.2
for (col, ttl, desc), (bx, by) in zip(boxes, grid):
    rect(s, bx, by, BW, BH, color=LIGHT)
    t(s, ttl,  bx+0.25, by+0.18, BW-0.4, 0.48, size=18, bold=True, color=col)
    t(s, desc, bx+0.25, by+0.72, BW-0.4, 1.2,  size=16, color=INK)

reason_slide(1, BLUE, "強 Agentic Loop",
    "不是 Autocomplete，而是 Iterative Problem Solver", [
    "工作方式是循環式：Gather Context → Take Action → Verify Results，不是一次性輸出",
    "更接近真實工程流程：真實寫 code 不是一次寫對，而是先理解、再改動、再驗證、再修正",
    "能持續根據結果修正方向：test failure、build error 都是下一輪推理的輸入，不是終點",
])
reason_slide(2, GREEN, "強 Context Management",
    "Loop 變長後，品質還能撐住", [
    "長任務的瓶頸常常不是模型智力，而是 context 汙染",
    "任務一長，就會累積大量雜訊（舊方案、失敗嘗試、大量 log），容易讓模型失焦",
    "探索、規劃、實作可以分開處理，避免主上下文被過多細節塞滿",
])
reason_slide(3, ORANGE, "強 Tools",
    "讓 Loop 每一步真正落地", [
    "Claude Code 不是只回答，而是真的能動手做：搜尋 codebase、編輯檔案、執行命令",
    "Tools 讓 reasoning 變成 action：沒有工具，loop 只能停在分析，無法實際執行",
    "Verify 之所以有價值，是因為真的能執行測試、lint、build，再根據結果修正",
])
reason_slide(4, PURPLE, "強 Model",
    "決定整個系統的推理上限與工具使用精度", [
    "前面有 loop、context management、tools，最後仍要靠 model 理解需求、規劃步驟、選對工具、解讀結果",
    "Tool calling 能力屬於 benchmark 領先群組（如 Berkeley Function Calling Leaderboard）",
    "依任務複雜度自動分配不同 model：探索任務用輕量模型（快 / 省成本），主推理用更強的模型",
])


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 – 案例帶入
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
t(s, "案例實際帶入", 0.65, 0.45, 11.5, 0.85, size=36, bold=True)
t(s, "折扣碼只能用一次，但用戶發現可以重複使用 — 四個 Reason 如何共同作用",
  0.65, 1.28, 11.5, 0.46, size=17, color=MUTED, italic=True)
hline(s, 1.82, x=0.65, w=12.03)

case_steps = [
    (ORANGE, "Tools",
     "grep 找到所有用到折扣碼的地方 → 讀各個驗證邏輯 → 確認哪裡沒有檢查是否已使用過",
     "沒有工具只能靠猜；有工具才能找出所有相關程式碼，不會有漏網之魚"),
    (BLUE, "Agentic Loop",
     "Round 1：修結帳 API 的驗證 → 還是能重複用　Round 2：修訂單建立的驗證 → 還是有漏　Round 3：在資料庫層加限制 → 全過",
     "每次測試失敗都是下一輪的輸入；autocomplete 給完第一個修法就停了"),
    (GREEN, "Context Mgmt",
     "同時追蹤結帳 API、訂單建立、折扣碼查詢邏輯，以及哪些地方已修、哪些待修",
     "多個地方同時進行，沒有管理容易遺漏，或把已修好的改回去"),
    (PURPLE, "Model",
     "判斷要在哪一層加驗證才能真正防住：API 層？service 層？還是資料庫層？推理出最根本的解法",
     "前面三步能 work，都建立在 model 每一步判斷正確：找對位置、選對修法、知道何時繼續迭代。任何一步推理錯，整個流程就跑偏"),
]
STEP_H = 1.28
cy = 1.95
for i, (col, reason, what, why) in enumerate(case_steps):
    rect(s, 0.65, cy, 12.03, STEP_H, color=LIGHT)
    t(s, f"0{i+1}", 0.82, cy+0.38, 0.55, 0.48, size=20, bold=True, color=col)
    t(s, what,        1.52, cy+0.10, 8.58, 0.52, size=15, bold=True, color=INK)
    t(s, "→ " + why,  1.52, cy+0.62, 8.58, 0.55, size=14, color=MUTED, italic=True)
    rect(s, 10.25, cy+0.39, 2.28, 0.50, color=col)
    t(s, reason, 10.27, cy+0.40, 2.24, 0.46,
      size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    cy += STEP_H + 0.07


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 – Transition
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.33, 7.5, color=LIGHT)
t(s, "Part 2",              0.85, 1.85, 11, 0.5,  size=18, bold=True, color=MUTED)
t(s, "機制深入",            0.85, 2.35, 11, 1.2,  size=58, bold=True, color=INK)
hline(s, 3.78, x=0.85, w=11.63)
t(s, "前面看了 Why — 現在來看 How",
  0.85, 3.92, 11, 0.52, size=22, color=MUTED, italic=True)
t(s, "Agentic Loop 的每一步，底層都是一個 Tool Call 的循環",
  0.85, 4.52, 11, 0.52, size=18, color=INK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 – ReAct Loop
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "01  Agent Architecture — ReAct Loop")

steps = ["使用者輸入", "LLM 推理", "Tool Call", "Tool Result", "LLM 推理", "…直到完成"]
cols  = [MUTED, BLUE, GREEN, ORANGE, BLUE, MUTED]
xpos  = [0.5, 2.45, 4.4, 6.35, 8.3, 10.25]
bw, bh = 1.82, 0.68
y_box = 1.32
for i, (step, col, x) in enumerate(zip(steps, cols, xpos)):
    rect(s, x, y_box, bw, bh, color=GRAY)
    t(s, step, x, y_box, bw, bh, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        t(s, "→", x+bw+0.01, y_box+0.12, 0.25, 0.4,
          size=18, color=MUTED, align=PP_ALIGN.CENTER)

hline(s, 2.2)
label(s, "重要特性", 0.55, 2.3)
bullets(s, [
    "單次 LLM 推理可在同一個 content 陣列中同時產生 text + 多個 tool_use",
    "多個 tool_use 批次發出 → 平行執行 → 所有結果一起回來 → 才進行下一次推理",
    "Loop 終止條件：LLM 推理產生的回應不含任何 tool_use，控制權回到使用者",
    "Background Bash（run_in_background）不阻塞推理，結果以非同步通知送達",
    "Subagent 有自己獨立的 agentic loop，不佔主 agent 的推理次數",
], 0.55, 2.68, 12.2, size=17, gap=0.50)


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 – 工具總覽 1/2
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "02  工具總覽 (1/2) — 核心工具", "★ = auto-allow，預設直接執行不詢問")

NX, NW = 0.50, 2.15
AX, AW = 2.70, 0.40
DX, DW = 3.18, 9.72
ROW_H  = 0.46

groups_1 = [
    (BLUE, "唯讀探索", [
        ("Read",      "★", "讀取本地檔案；支援圖片、PDF、ipynb；25k token 上限"),
        ("Glob",      "★", "用 glob pattern 搜尋符合條件的檔案路徑"),
        ("Grep",      "★", "直接 spawn ripgrep 搜尋檔案內容（bypass shell，無注入風險）"),
        ("WebFetch",  " ", "抓取 URL 網頁內容，HTML 自動轉 markdown"),
        ("WebSearch", " ", "搜尋網路，取得 training data 以外的最新資訊"),
    ]),
    (GREEN, "寫入", [
        ("Edit",         " ", "精確 old→new 替換；old_string 必須唯一；只傳 diff，token 最少"),
        ("Write",        " ", "覆寫整個檔案；適合新建或完整重寫；自動產生 diff"),
        ("NotebookEdit", " ", "編輯 .ipynb cell（必須用此；直接 Edit 會報錯）"),
    ]),
    (ORANGE, "執行", [
        ("Bash", " ", "執行任意 shell 命令；支援 run_in_background；需注意 shell injection"),
    ]),
]

cy = 1.22
for col, grp_name, tools in groups_1:
    rect(s, 0.50, cy, 12.33, 0.33, color=col)
    t(s, grp_name, NX+0.1, cy+0.06, 3, 0.26, size=13, bold=True, color=WHITE)
    cy += 0.35
    for name, star, desc in tools:
        hline(s, cy, color=GRAY)
        t(s, name, NX+0.1,  cy+0.07, NW,  ROW_H-0.1, size=14, bold=True, color=col)
        t(s, star, AX+0.05, cy+0.07, AW,  ROW_H-0.1, size=14, bold=True, color=BLUE)
        t(s, desc, DX+0.05, cy+0.07, DW,  ROW_H-0.1, size=13, color=INK)
        cy += ROW_H
hline(s, cy, color=GRAY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 – 工具總覽 2/2
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "02  工具總覽 (2/2) — 系統與擴充工具")

NX2, NW2 = 0.50, 3.60
DX2, DW2 = 4.18, 8.72
ROW_H2   = 0.42

groups_2 = [
    (MUTED, "Meta / 系統", [
        ("ToolSearch",       "★", "將 deferred 工具 schema 載入 context，才能呼叫它"),
        ("Agent",            " ", "啟動有獨立 context 的 subagent（可用不同 model）"),
        ("Skill",            " ", "呼叫 .claude/skills/ 中定義的自訂指令集"),
        ("AskUserQuestion",  " ", "主動向使用者提問（唯一能讓 Claude 暫停等待輸入的工具）"),
    ]),
    (RED, "任務管理", [
        ("TaskCreate / TaskUpdate / TaskStop", " ", "建立、更新、停止任務；TaskStop 只能停背景 Bash 任務"),
        ("TaskGet / TaskList / TaskOutput",    "★", "查詢任務狀態與輸出"),
        ("CronCreate / CronDelete / CronList", " ", "建立 session-only 定時排程（Claude 退出後消失）"),
    ]),
    (PURPLE, "規劃 / 隔離", [
        ("EnterPlanMode / ExitPlanMode", " ", "進入規劃模式；計畫存至 ~/.claude/plans/"),
        ("EnterWorktree",                " ", "在 .claude/worktrees/ 建立隔離的 git worktree 環境"),
    ]),
    (TEAL, "IDE (MCP)", [
        ("mcp__ide__getDiagnostics", " ", "取得 IDE lint / type 錯誤診斷（VS Code / JetBrains）"),
        ("mcp__ide__executeCode",    " ", "在 IDE 環境執行程式碼片段"),
    ]),
]

cy = 1.22
for col, grp_name, tools in groups_2:
    rect(s, 0.50, cy, 12.33, 0.33, color=col)
    t(s, grp_name, NX2+0.1, cy+0.06, 4, 0.26, size=13, bold=True, color=WHITE)
    cy += 0.35
    for name, star, desc in tools:
        hline(s, cy, color=GRAY)
        t(s, name, NX2+0.1, cy+0.06, NW2,  ROW_H2, size=13, bold=True, color=col)
        t(s, desc, DX2+0.1, cy+0.06, DW2,  ROW_H2, size=13, color=INK)
        cy += ROW_H2
hline(s, cy, color=GRAY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 – Read vs Grep vs Glob
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "03  工具深入 — Read / Grep / Glob",
             "同樣 auto-allow 且唯讀，但設計定位不同，不要用 Bash 代替")

col_x = [0.50, 3.05, 6.50, 9.95]
col_w = [2.48, 3.38, 3.38, 3.25]
rect(s, 0.50, 1.28, 12.33, 0.42, color=GRAY)
for h, x, w, c in zip(["", "Read", "Grep", "Glob"],
                       col_x, col_w,
                       [MUTED, BLUE, GREEN, ORANGE]):
    t(s, h, x+0.1, 1.31, w, 0.36, size=16, bold=True, color=c)

rows = [
    ("用途",      "讀取檔案完整內容",           "搜尋哪些位置含特定字串",      "搜尋哪些路徑符合條件"),
    ("輸入",      "檔案路徑",                   "regex pattern + 目錄",        "glob pattern"),
    ("特殊能力",  "圖片/PDF/ipynb 特殊解析",    "直接 spawn ripgrep，不走 shell","結果依修改時間排序"),
    ("token 限制","25,000 tokens 硬上限",       "無",                          "無"),
    ("典型用法",  "看已知檔案的內容",           "找出哪個檔案含這段邏輯",      "找出有哪些 .ts 檔案"),
]
cy = 1.73
for row in rows:
    hline(s, cy, color=GRAY)
    colors = [MUTED, BLUE, GREEN, ORANGE]
    for val, x, w, c in zip(row, col_x, col_w, colors):
        bold = (c == MUTED)
        t(s, val, x+0.1, cy+0.08, w-0.15, 0.46, size=14, bold=bold, color=c)
    cy += 0.54
hline(s, cy, color=GRAY)

label(s, "實際使用順序範例 — 三種工具如何接力", 0.55, cy+0.14, color=BLUE, size=13)
t(s, "情境：找出折扣碼驗證邏輯在哪裡，再完整閱讀",
  0.55, cy+0.48, 12.2, 0.30, size=14, bold=True, color=MUTED)
rect(s, 0.55, cy+0.82, 12.2, 1.38, color=LIGHT)
t(s, "① Glob(\"**/*.ts\")  →  列出所有 TypeScript 檔案路徑清單，確認搜尋範圍",
  0.73, cy+0.92, 11.8, 0.28, size=14, color=INK)
t(s, "② Grep(\"validateDiscount\", \"**/*.ts\")  →  找出含該函式的檔案（例：discount.service.ts）",
  0.73, cy+1.20, 11.8, 0.28, size=14, color=INK)
t(s, "③ Read(\"src/services/discount.service.ts\")  →  讀取完整內容，理解驗證邏輯細節",
  0.73, cy+1.48, 11.8, 0.28, size=14, color=INK)
t(s, "設計目的：三種工具皆 auto-allow（無確認中斷），從路徑定位到內容閱讀一氣呵成，比 Bash 更安全、更省確認次數",
  0.55, cy+2.28, 12.2, 0.30, size=13, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 – Edit vs Write
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "03  工具深入 — Edit vs Write")

col_x2 = [0.50, 3.55, 8.10]
col_w2 = [2.98, 4.48, 4.82]
rect(s, 0.50, 1.28, 12.33, 0.42, color=GRAY)
for h, x, w, c in zip(["面向", "Edit", "Write"],
                       col_x2, col_w2, [MUTED, GREEN, BLUE]):
    t(s, h, x+0.1, 1.31, w, 0.36, size=16, bold=True, color=c)

rows2 = [
    ("操作方式",   "精確字串替換（old_string → new_string）",  "覆寫整個檔案內容"),
    ("適用情境",   "修改現有檔案的一部分",                    "新建檔案、或整個檔案都要換"),
    ("唯一性限制", "old_string 必須在檔案中唯一，否則報錯",   "無"),
    ("token 消耗", "只傳 diff（少）",                         "傳送整個檔案（多）"),
    ("diff 顯示",  "自動計算 structuredPatch + gitDiff",       "同左"),
    ("replace_all","支援：一次替換所有匹配",                  "不適用"),
]
cy = 1.73
for dim, e, w in rows2:
    hline(s, cy, color=GRAY)
    t(s, dim, col_x2[0]+0.1, cy+0.08, col_w2[0]-0.15, 0.44, size=14, bold=True, color=MUTED)
    t(s, e,   col_x2[1]+0.1, cy+0.08, col_w2[1]-0.15, 0.44, size=14, color=GREEN)
    t(s, w,   col_x2[2]+0.1, cy+0.08, col_w2[2]-0.15, 0.44, size=14, color=BLUE)
    cy += 0.52
hline(s, cy, color=GRAY)

label(s, "使用原則", 0.55, cy+0.14, size=13)
t(s, "日常修改首選 Edit（token 省、定位精準）；只有新建檔案或整個內容都要換時才用 Write。",
  0.55, cy+0.50, 12.2, 0.42, size=16, color=INK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 15 – Bash
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "03  工具深入 — Bash",
             "有專用工具時優先用專用工具；只有需要『執行』時才用 Bash")

# Decision table
label(s, "專用工具 vs Bash — 怎麼選？", 0.55, 1.12, color=BLUE)
THEAD_Y = 1.48
col_xb = [0.50, 3.20, 6.40, 10.00]
col_wb = [2.62, 3.12, 3.52, 3.20]
rect(s, 0.50, THEAD_Y, 12.33, 0.36, color=GRAY)
for h, x, w, c in zip(["需要做的事", "用這個工具", "為何不用 Bash", "Auto-allow?"],
                       col_xb, col_wb, [MUTED, GREEN, MUTED, MUTED]):
    t(s, h, x+0.10, THEAD_Y+0.06, w, 0.28, size=13, bold=True, color=c)

dec_rows = [
    ("搜尋哪個檔案含某字串",      "Grep",       "直接 spawn ripgrep，更快且無 injection 風險",   "✓"),
    ("找出哪些路徑符合條件",      "Glob",       "結果依修改時間排序，且不走 shell",               "✓"),
    ("讀取檔案內容",              "Read",       "支援 PDF / 圖片 / ipynb，有 token 保護上限",     "✓"),
    ("修改檔案一部分",            "Edit",       "只傳 diff，token 省；old_string 定位精準",        "✗（需確認）"),
    ("版控 / 套件管理 / 跑測試",  "Bash ✓",     "需要真實 shell 環境，專用工具無法執行",          "✗（需確認）"),
]
cy = THEAD_Y + 0.36
for task, tool, reason, allow in dec_rows:
    hline(s, cy, color=GRAY)
    is_bash = tool.startswith("Bash")
    row_col = LIGHT if not is_bash else RGBColor(0xFF, 0xF3, 0xE6)
    rect(s, 0.50, cy, 12.33, 0.44, color=row_col)
    t(s, task,   col_xb[0]+0.10, cy+0.08, col_wb[0], 0.34, size=13, color=INK)
    t(s, tool,   col_xb[1]+0.10, cy+0.08, col_wb[1], 0.34, size=13,
      bold=is_bash, color=ORANGE if is_bash else GREEN)
    t(s, reason, col_xb[2]+0.10, cy+0.08, col_wb[2], 0.34, size=12, color=MUTED)
    t(s, allow,  col_xb[3]+0.10, cy+0.08, col_wb[3], 0.34, size=13,
      bold=True, color=GREEN if allow == "✓" else RED)
    cy += 0.44
hline(s, cy, color=GRAY)

hline(s, cy + 0.12)
label(s, "Bash 執行特性", 0.55, cy+0.24)

cy2 = bullets(s, [
    "起完整 bash shell：支援 &&、pipe、環境變數展開，與手動在終端機執行完全相同",
    "run_in_background：命令移到背景執行，Claude 繼續推理；完成後才收到 <task-notification> 非同步通知",
    "timeout：預設 2 分鐘；長任務可設最長 10 分鐘（如 npm install）",
    "exit code：非 0 時 Claude 自動讀取 stderr，分析錯誤並嘗試修正",
], 0.55, cy+0.62, 12.2, size=15, gap=0.40)

hline(s, cy2+0.10)
label(s, "安全考量", 0.55, cy2+0.22, color=RED, size=13)
bullets(s, [
    "命令字串經過 shell 解析，有 injection 風險；Read/Edit 由 framework 直接組裝參數，沒有此問題",
    "複合命令（&&, ||, ;, |）永遠不匹配 allow 規則 → 防止 'git status && rm -rf /' 越獄",
], 0.55, cy2+0.58, 12.2, size=15, gap=0.40)


# ══════════════════════════════════════════════════════════════════
# SLIDE 16 – ToolSearch
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "03  工具深入 — ToolSearch",
             "唯一預設載入的 meta 工具：讓其他工具的 schema 進入 context")

t(s, "工具 schema 全部預載會消耗大量 token，設計上分為兩種模式，由 ToolSearch 負責按需載入：",
  0.55, 1.18, 12.2, 0.34, size=15, color=INK)

BX1, BX2, BW = 0.55, 6.90, 6.00
BY = 1.56
rect(s, BX1, BY, BW, 0.36, color=BLUE)
t(s, "Pre-loaded（9）— session 開始即載入，零延遲",
  BX1+0.12, BY+0.06, BW-0.2, 0.28, size=13, bold=True, color=WHITE)
rect(s, BX2, BY, BW, 0.36, color=MUTED)
t(s, "Deferred（17）— 需先 ToolSearch 載入，+2-3 秒",
  BX2+0.12, BY+0.06, BW-0.2, 0.28, size=13, bold=True, color=WHITE)

preloaded = ["Read", "Glob", "Grep", "Edit", "Write",
             "Bash", "ToolSearch", "Agent", "Skill"]
deferred_a = ["AskUserQuestion", "WebFetch", "WebSearch",
              "NotebookEdit", "TaskCreate", "TaskUpdate",
              "TaskStop", "TaskGet", "TaskList"]
deferred_b = ["TaskOutput", "CronCreate", "CronDelete",
              "CronList", "EnterPlanMode", "ExitPlanMode",
              "EnterWorktree", "ExitWorktree"]

cy_l = BY + 0.42
for name in preloaded:
    t(s, "• " + name, BX1+0.15, cy_l, 5.6, 0.26, size=13, color=INK)
    cy_l += 0.27

cy_r = BY + 0.42
for name in deferred_a:
    t(s, "• " + name, BX2+0.12, cy_r, 2.7, 0.26, size=13, color=INK)
    cy_r += 0.27

cy_r2 = BY + 0.42
for name in deferred_b:
    t(s, "• " + name, BX2+3.12, cy_r2, 2.7, 0.26, size=13, color=INK)
    cy_r2 += 0.27

list_bottom = BY + 0.42 + len(preloaded) * 0.27   # ≈ 4.41

hline(s, list_bottom + 0.12)
label(s, "使用流程範例：呼叫 NotebookEdit（Deferred 工具）",
  0.55, list_bottom + 0.22, color=ORANGE, size=13)
rect(s, 0.55, list_bottom + 0.60, 12.2, 1.72, color=LIGHT)
fy = list_bottom + 0.72
t(s, "❶  Claude 嘗試呼叫 NotebookEdit，但它在 <available-deferred-tools>（有名稱但無 schema，無法直接呼叫）",
  0.73, fy, 11.8, 0.34, size=13, color=INK)
t(s, '❷  ToolSearch("select:NotebookEdit")  →  +1 round-trip，schema 載入 context（約 2-3 秒）',
  0.73, fy+0.42, 11.8, 0.34, size=13, color=INK)
t(s, "❸  NotebookEdit 現在可正常呼叫來編輯 .ipynb cell；後續同 session 不需再次 ToolSearch",
  0.73, fy+0.84, 11.8, 0.34, size=13, color=INK)
t(s, "使用方式：select:ToolName（精確）或關鍵字搜尋如 'notebook jupyter'（不確定名稱時）",
  0.73, fy+1.28, 11.8, 0.30, size=12, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 17 – Agent
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  工具深入 — Agent Tool（Subagent）")

label(s, "5 種 subagent_type", 0.50, 1.15)
# (name, tools, role, example)
types = [
    ("general-purpose",  "全部工具（*）",
     "通用型，處理複雜多步驟任務",
     "例：獨立重構 auth 模組，讀寫多個檔案、跑測試、循環修錯誤"),
    ("Explore",          "除 Edit、Write 等寫入工具",
     "快速掃描 codebase（quick / medium / very thorough 三種深度）",
     "例：「了解 src/payment/ 的模組架構與依賴關係」"),
    ("Plan",             "同 Explore",
     "回傳 step-by-step 實作計畫，不直接改程式碼",
     "例：「規劃把 REST API 改成 GraphQL 的步驟與影響範圍」"),
    ("claude-code-guide","Glob、Grep、Read、WebFetch、WebSearch",
     "回答 Claude Code / SDK / API 使用問題；支援 resume 繼續上次對話",
     "例：「Hooks 的 PreToolUse 怎麼設定阻斷？」"),
    ("statusline-setup", "Read、Edit",
     "設定 Claude Code status line 顯示內容",
     "例：「在 status bar 顯示目前使用的 model 名稱」"),
]
col_x3 = [0.50, 3.30, 7.00]
col_w3 = [2.72, 3.60, 5.60]
rect(s, 0.50, 1.52, 12.33, 0.38, color=GRAY)
for h, x, w in zip(["subagent_type", "可用工具", "用途 & 使用情境"], col_x3, col_w3):
    t(s, h, x+0.08, 1.54, w, 0.32, size=13, bold=True, color=MUTED)
cy = 1.93
for name, tools_str, role, example in types:
    hline(s, cy, color=GRAY)
    t(s, name,      col_x3[0]+0.08, cy+0.06, col_w3[0], 0.30, size=13, bold=True, color=BLUE)
    t(s, tools_str, col_x3[1]+0.08, cy+0.06, col_w3[1], 0.52, size=12, color=INK)
    t(s, role,      col_x3[2]+0.08, cy+0.06, col_w3[2], 0.28, size=13, bold=False, color=INK)
    t(s, example,   col_x3[2]+0.08, cy+0.34, col_w3[2], 0.26, size=12, color=MUTED, italic=True)
    cy += 0.66
hline(s, cy, color=GRAY)

label(s, "關鍵特性", 0.50, cy+0.12)
bullets(s, [
    "獨立 context：每個 subagent 有自己的 context window，不繼承主 agent 對話歷史 → 防止主 context 被細節塞滿，這是 Context Management 的核心實作",
    "不同 model：Explore subagent 實測使用 claude-haiku（快 / 省 token），主 agent 使用 claude-sonnet → 依任務複雜度自動分配 model",
    "agentId resume：帶入 resume 參數可繼續同一個 subagent 的 context，不用重新建立",
], 0.50, cy+0.50, 12.2, size=15, gap=0.46)


# ══════════════════════════════════════════════════════════════════
# SLIDE 18 – Skill
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  工具深入 — Skill 系統")

label(s, "Skill 是什麼", 0.50, 1.15)
bullets(s, [
    "放在 .claude/skills/<name>/SKILL.md 的 Markdown 指令檔，告訴 Claude 如何完成特定任務",
    "可附帶腳本，在展開給 LLM 之前先預處理資料（! 前綴命令，LLM 看不到原始命令）",
], 0.50, 1.52, 12.2, size=17, gap=0.50)

hline(s, 2.58)
label(s, "兩種觸發方式", 0.50, 2.68)

rect(s, 0.50, 3.08, 5.90, 0.38, color=GRAY)
t(s, "Trigger A  /skill-name  （使用者輸入）",
  0.62, 3.11, 5.65, 0.32, size=15, bold=True, color=BLUE)
bullets(s, [
    "CLI 在框架層讀取 SKILL.md，展開為 isMeta:true 訊息注入 context",
    "Claude 不需呼叫 Skill tool，直接看到展開後的指令",
], 0.62, 3.48, 5.65, size=16, gap=0.46)

rect(s, 6.90, 3.08, 5.93, 0.38, color=GRAY)
t(s, "Trigger B  Skill tool  （Claude 主動呼叫）",
  7.02, 3.11, 5.70, 0.32, size=15, bold=True, color=PURPLE)
bullets(s, [
    'Claude 先 ToolSearch 載入 Skill schema，再呼叫 Skill("name", args)',
    "Framework 處理 SKILL.md → tool_result 回傳給 LLM",
], 7.02, 3.48, 5.70, size=16, gap=0.46)

hline(s, 4.46)
label(s, "Skill allowed-tools vs Agent tools 的本質差異", 0.50, 4.54, color=RED, size=13)
bullets(s, [
    "Skill allowed-tools = 軟性：LLM 仍可呼叫其他工具（需使用者確認），context 共享",
    "Agent tools = 硬性：API 層隔離，LLM 根本看不到白名單以外的工具 schema",
], 0.50, 4.90, 12.2, size=16, gap=0.50)


# ══════════════════════════════════════════════════════════════════
# SLIDE 19 – Task Tools
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  工具深入 — Task 任務管理工具",
             "追蹤 Bash run_in_background 產生的背景任務；★ = auto-allow")

t(s, "Bash 呼叫 run_in_background=True 時，命令立即移往背景執行並返回 task_id，Task 工具負責追蹤這些背景任務。",
  0.50, 1.15, 12.2, 0.36, size=15, color=INK)

task_tools = [
    ("TaskCreate",  " ", "建立任務記錄（通常由 Bash run_in_background 自動觸發）"),
    ("TaskUpdate",  " ", "更新任務的 metadata 或描述"),
    ("TaskStop",    " ", "強制停止一個正在執行的背景任務"),
    ("TaskGet",     "★", "取得單一任務的目前狀態（running / completed / failed）"),
    ("TaskList",    "★", "列出 session 中所有任務及其狀態"),
    ("TaskOutput",  "★", "取得任務的完整 stdout / stderr 輸出內容"),
    ("CronCreate",  " ", "建立 session-only 定時排程（Claude 退出後消失）"),
    ("CronDelete",  " ", "刪除定時排程"),
    ("CronList",    "★", "列出所有定時排程"),
]
NXT2, NXA2, NXD2 = 0.50, 2.68, 3.18
NWT2, NWA2, NWD2 = 2.10, 0.42, 9.68
ROW_HT = 0.38

rect(s, 0.50, 1.58, 12.33, 0.34, color=GRAY)
t(s, "工具名稱", NXT2+0.10, 1.61, NWT2, 0.26, size=13, bold=True, color=MUTED)
t(s, "★",       NXA2+0.05, 1.61, NWA2, 0.26, size=13, bold=True, color=BLUE)
t(s, "功能說明", NXD2+0.05, 1.61, NWD2, 0.26, size=13, bold=True, color=MUTED)

cy = 1.94
for name, star, desc in task_tools:
    hline(s, cy, color=GRAY)
    t(s, name, NXT2+0.10, cy+0.06, NWT2,  ROW_HT-0.06, size=13, bold=True, color=TEAL)
    t(s, star, NXA2+0.05, cy+0.06, NWA2,  ROW_HT-0.06, size=13, bold=True, color=BLUE)
    t(s, desc, NXD2+0.05, cy+0.06, NWD2,  ROW_HT-0.06, size=13, color=INK)
    cy += ROW_HT
hline(s, cy, color=GRAY)

hline(s, cy + 0.14)
label(s, "使用流程範例：背景跑測試，完成後取結果", 0.50, cy+0.26, color=TEAL, size=13)
rect(s, 0.50, cy+0.64, 12.2, 1.08, color=LIGHT)
fy = cy + 0.74
t(s, "❶  Bash(\"npm test --all\", run_in_background=True)  →  立即返回 task_id，Claude 不等待，繼續推理其他工作",
  0.68, fy,        11.8, 0.28, size=13, color=INK)
t(s, "❷  收到 <task-notification: task_id completed>  →  Claude 得知任務結束，可以去取結果了",
  0.68, fy+0.30,   11.8, 0.28, size=13, color=INK)
t(s, "❸  TaskOutput(task_id)  →  取得完整 test 輸出，分析哪些失敗，進行下一輪修正",
  0.68, fy+0.60,   11.8, 0.28, size=13, color=INK)
t(s, "CronCreate 則用於定時排程，例如：每 10 分鐘自動執行一次 lint check，session 期間持續運作",
  0.50, cy+1.80,   12.2, 0.28, size=13, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 20 – Hooks
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "05  系統機制 — Hooks（事件驅動擴充）")

label(s, "關鍵 Hook 事件（共 21 種，選幾個最常用的）", 0.50, 1.15)

key_hooks = [
    (GREEN,  "PreToolUse",       "工具呼叫前",
     "exit 2 可阻斷工具呼叫。例：偵測到 Bash rm 命令時自動阻止，防止意外刪除檔案"),
    (ORANGE, "PostToolUse",      "工具執行後",
     "例：Edit 存檔後自動跑 eslint，確保每次寫入都符合程式碼規範，不需手動執行"),
    (BLUE,   "UserPromptSubmit", "使用者送出訊息前",
     "例：自動在每個 prompt 前注入 git status，讓 Claude 始終知道目前工作狀態"),
    (MUTED,  "SessionStart",     "對話開始時",
     "例：載入當前 sprint 的 ticket 列表，讓 Claude 自動了解本次任務範圍"),
    (PURPLE, "Stop",             "Claude 完成工作後",
     "例：自動 git push 並發 Slack 通知，或將工作摘要寫入 daily log"),
]
ROW_H = 0.64
cy = 1.52
for col, event, timing, desc in key_hooks:
    rect(s, 0.50, cy, 12.33, ROW_H, color=LIGHT)
    t(s, event,  0.62, cy+0.06, 2.20, 0.28, size=14, bold=True, color=col)
    t(s, timing, 0.62, cy+0.36, 2.20, 0.24, size=11, color=MUTED, italic=True)
    t(s, desc,   3.00, cy+0.10, 9.70, 0.46, size=13, color=INK)
    cy += ROW_H + 0.04

hline(s, cy + 0.08)
label(s, "Hook 的四種執行類型 — 決定 hook 觸發後「怎麼做」", 0.50, cy+0.20)

# 2x2 grid of hook types
BX1H, BX2H = 0.50, 6.72
BWH, BHH = 6.08, 0.72
BY1H = cy + 0.58
BY2H = BY1H + BHH + 0.06

type_data = [
    (BX1H, BY1H, GREEN,  "command — 執行 shell 命令",
     "最常用。exit 2 可阻斷；async 背景執行。",
     "例：PostToolUse → 存檔後自動跑 eslint"),
    (BX2H, BY1H, PURPLE, "prompt — 讓另一個 LLM 評估",
     "用輕量 model（Haiku）判斷是否允許，回傳 {ok: true/false}。",
     "例：PreToolUse → 判斷 Bash 命令是否有危險"),
    (BX1H, BY2H, BLUE,   "http — POST 到指定 URL",
     "呼叫 webhook，不需起額外程序，適合通知類場景。",
     "例：Stop → 發 Slack 通知「Claude 完成了」"),
    (BX2H, BY2H, ORANGE, "agent — 起 subagent 做複雜驗證",
     "最強大；subagent 有獨立 context 與工具，可多輪推理。",
     "例：PreToolUse → subagent 確認程式碼修改不破壞 API 契約"),
]
for bx, by, col, title, body, example in type_data:
    rect(s, bx, by, BWH, BHH, color=LIGHT)
    t(s, title,   bx+0.12, by+0.06, BWH-0.2, 0.26, size=13, bold=True, color=col)
    t(s, body,    bx+0.12, by+0.34, BWH-0.2, 0.22, size=12, color=INK)
    t(s, example, bx+0.12, by+0.50, BWH-0.2, 0.20, size=11, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 20 – Permission
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "05  系統機制 — Permission 系統")

# ── Section A: 模型本身的推論習慣（軟性）────────────────────────
rect(s, 0.50, 1.15, 12.33, 0.38, color=RED)
t(s, "Claude 模型的推論習慣（軟性限制）", 0.65, 1.18, 11, 0.30, size=15, bold=True, color=WHITE)
t(s, "這是 Claude 訓練出來的謹慎行為，不是程式碼規則。無法透過任何設定改變，因為改的是 framework，不是模型本身。",
  0.65, 1.60, 12.0, 0.32, size=14, color=INK)
soft_items = [
    "操作超出使用者明確要求的範圍  →  傾向先問再做（例：被要求修 A，發現連 B 也有問題，會先確認）",
    "發現非預期的檔案、分支、或狀態  →  傾向先回報再繼續（例：發現有未提交的修改）",
    "任務指令模糊、有多種解讀方式  →  傾向暫停確認而不是自己猜",
]
bullets(s, soft_items, 0.65, 1.96, 12.05, size=14, gap=0.44)

# ── Section B: Tool 硬限制規則（由 framework 執行）──────────────
rect(s, 0.50, 3.38, 12.33, 0.38, color=BLUE)
t(s, "Tool 硬限制規則（由 framework 強制執行）", 0.65, 3.41, 11, 0.30, size=15, bold=True, color=WHITE)

hard_rules = [
    (GREEN,  "Auto-allow",
     "Read、Glob、Grep、ToolSearch 等唯讀工具預設直接放行，不詢問",
     "目的：讓 Claude 自由探索，維持 Agentic Loop 流暢"),
    (ORANGE, "Plan Mode",
     "EnterPlanMode 後，Edit / Write / Bash 等所有寫入 / 執行工具一律禁止",
     "目的：規劃階段保證不會意外修改程式碼"),
    (BLUE,   "Settings.json allow / deny",
     'permissions.allow 預先核准特定工具：「Edit」、「Bash(git:*)」、「WebFetch(domain:...)」',
     "符合規則的命令直接放行；deny 則永久封鎖特定工具"),
]
cy = 3.82
for col, title, line1, line2 in hard_rules:
    rect(s, 0.50, cy, 12.33, 1.02, color=LIGHT)
    rect(s, 0.50, cy, 0.20, 1.02, color=col)
    t(s, title, 0.82, cy+0.08, 3.0,  0.28, size=14, bold=True, color=col)
    t(s, line1, 0.82, cy+0.40, 11.8, 0.26, size=13, color=INK)
    t(s, line2, 0.82, cy+0.70, 11.8, 0.26, size=13, color=MUTED, italic=True)
    cy += 1.06

t(s, "Bash allow 格式：prefix Bash(git:*) 匹配任何 git 開頭命令；複合命令（&&, ||）永遠不匹配，防越獄",
  0.50, cy+0.10, 12.2, 0.28, size=12, color=MUTED, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 22 – 資訊來源
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "06  資訊來源 — 本簡報的機制如何驗證",
             "三種來源各有優先順序：JSONL > 官方文件 > 原始碼")

sources = [
    (BLUE,   "Session JSONL（最優先，最可信）",
     "Claude Code 在本機記錄每次對話的完整 JSONL，包含所有 tool call、推理內容、token 用量",
     "記錄的是實際發生的行為，不是設計意圖。JSONL 與文件矛盾時，以 JSONL 為準。",
     "路徑：C:\\Users\\<user>\\.claude\\projects\\<project>\\<session-id>.jsonl"),
    (GREEN,  "官方文件（次優先）",
     "code.claude.com/docs — 說明設計意圖與公開行為，快速查閱各功能用途",
     "適合確認功能設計，但不一定涵蓋所有實作細節（例如 binary 層級的行為）",
     "索引：https://code.claude.com/docs/llms.txt  （base URL：code.claude.com/docs/en/）"),
    (ORANGE, "Claude Code 原始碼（最後手段）",
     "Claude Code 是單一執行檔（Node.js 打包），內嵌 minified JS bundle，沒有加密",
     "可用 Python regex 或 strings + grep 搜尋邏輯，適合找文件沒寫到的細節實作",
     "執行檔：C:\\Users\\<user>\\.local\\share\\claude\\versions\\<version>（約 237 MB）"),
]

cy = 1.28
for col, title, desc, note, path in sources:
    rect(s, 0.50, cy, 12.33, 1.58, color=LIGHT)
    rect(s, 0.50, cy, 0.22, 1.58, color=col)
    t(s, title, 0.84, cy+0.10, 11.8, 0.30, size=16, bold=True, color=col)
    t(s, desc,  0.84, cy+0.46, 11.8, 0.30, size=14, color=INK)
    t(s, note,  0.84, cy+0.80, 11.8, 0.28, size=13, color=MUTED, italic=True)
    t(s, path,  0.84, cy+1.12, 11.8, 0.28, size=12, color=MUTED)
    cy += 1.64


# ══════════════════════════════════════════════════════════════════
# SLIDE 23 – Takeaways
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "06  Takeaways")

takes = [
    (BLUE,   "Agentic Loop",
     "批次 tool call 平行執行，不含 tool_use 才停止。Subagent 有獨立 loop，不佔主 agent 的推理次數。"),
    (GREEN,  "工具設計哲學",
     "讀取工具預載且 auto-allow → 零中斷探索。Edit 只傳 diff → token 最少。每個工具針對最少確認次數最佳化。"),
    (ORANGE, "Pre-loaded vs Deferred",
     "9 個預載（零延遲）vs 16 個延遲載入（省 token，透過 ToolSearch，+2-3 秒）。根據使用頻率權衡 latency 與 token budget。"),
    (PURPLE, "Skill & Agent 的隔離差異",
     "Skill allowed-tools 是軟性（context 共享）。Agent tools 是硬性 API 層隔離（LLM 根本看不到其他工具 schema）。"),
    (RED,    "Permission 安全設計",
     "複合命令永遠不匹配 prefix 規則（binary 逆向確認）。Layer 2 判斷性確認無法用設定關閉。"),
    (MUTED,  "Strong Model 貫穿始終",
     "前面每一個機制能正確運作，背後都是 model 在每一步做出正確判斷。工具再多，model 推理錯誤，整個系統就跑偏。"),
]
cy = 1.28
for col, title, body in takes:
    hline(s, cy)
    t(s, title, 0.55, cy+0.10, 2.60, 0.48, size=17, bold=True, color=col)
    t(s, body,  3.25, cy+0.10, 9.55, 0.55, size=16, color=INK)
    cy += 0.82
hline(s, cy)


# ── Save ──────────────────────────────────────────────────────────
out = "D:/project/test_claude_skill/claude_code_full.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
