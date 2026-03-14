from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (white background, clean) ─────────────────────────
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1A, 0x22, 0x2E)   # primary text
MUTED  = RGBColor(0x55, 0x65, 0x7A)   # secondary text
BLUE   = RGBColor(0x1D, 0x6F, 0xC8)   # accent
GRAY   = RGBColor(0xE2, 0xE8, 0xF0)   # divider / light bg
GREEN  = RGBColor(0x05, 0x7A, 0x50)
ORANGE = RGBColor(0xB4, 0x5A, 0x00)
RED    = RGBColor(0xAA, 0x18, 0x18)
PURPLE = RGBColor(0x5B, 0x21, 0xB6)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Core helpers ───────────────────────────────────────────────

def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def hline(slide, y, x=0.5, w=12.33, color=GRAY, h=0.045):
    bar = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def rect(slide, x, y, w, h, color=GRAY):
    r = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    return r

def t(slide, text, x, y, w, h=0.55,
      size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def slide_header(slide, title, subtitle=None):
    bg(slide)
    hline(slide, 0.0, x=0.0, w=13.33, color=BLUE, h=0.10)
    t(slide, title, 0.55, 0.18, 12.2, 0.7, size=28, bold=True, color=INK)
    hline(slide, 1.0, color=GRAY)
    if subtitle:
        t(slide, subtitle, 0.55, 1.08, 12, 0.45, size=15, color=MUTED, italic=True)

def bullets(slide, items, x, y, w, size=18, gap=0.46, color=INK, indent="  "):
    """Render a list of strings as bullet lines. Prefix '>>' = sub-bullet."""
    cy = y
    for item in items:
        is_sub = item.startswith(">>")
        text = ("    • " + item[2:]) if is_sub else ("• " + item)
        c = MUTED if is_sub else color
        sz = size - 2 if is_sub else size
        t(slide, text, x, cy, w, gap, size=sz, color=c)
        cy += gap
    return cy

def label(slide, text, x, y, color=BLUE, size=14):
    """Small all-caps section label."""
    t(slide, text.upper(), x, y, 4, 0.35, size=size, bold=True, color=color)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, 13.33, 4.2, color=BLUE)
t(s, "Claude Code", 0.8, 0.7, 11.5, 1.4, size=60, bold=True, color=WHITE)
t(s, "Internal Mechanism Deep Dive", 0.8, 2.15, 11, 0.75, size=26, color=RGBColor(0xBB,0xD6,0xF8))
hline(s, 4.22, x=0, w=13.33, color=BLUE, h=0.06)
t(s, "Hsiang-Wen", 0.8, 4.65, 5, 0.5, size=20, color=INK)
t(s, "2026", 0.8, 5.22, 3, 0.4, size=16, color=MUTED)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 – Agenda
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "Agenda")

items = [
    ("01", "ReAct Loop — Claude Code 的基本運作方式"),
    ("02", "工具總覽 — 依功能分群"),
    ("03", "Claude Code 怎麼寫 Code — 流程與每步驟的設計優勢"),
    ("04", "Tool 特性深入說明"),
    ("05", "Permission 系統"),
    ("06", "Takeaways"),
]
cy = 1.35
for num, title in items:
    hline(s, cy, color=GRAY)
    t(s, num, 0.55, cy+0.1, 0.75, 0.5, size=20, bold=True, color=BLUE)
    t(s, title, 1.4, cy+0.1, 11, 0.5, size=20, color=INK)
    cy += 0.88


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 – ReAct Loop
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "01  ReAct Loop — Claude Code 的基本運作方式")

# Flow row
steps = ["使用者輸入", "LLM 推理", "Tool Call", "Tool Result", "LLM 推理", "…直到完成"]
cols  = [MUTED,        BLUE,        GREEN,       ORANGE,        BLUE,        MUTED]
xpos  = [0.5, 2.45, 4.4, 6.35, 8.3, 10.25]
bw, bh = 1.82, 0.68
y_box = 1.32
for i, (step, col, x) in enumerate(zip(steps, cols, xpos)):
    rect(s, x, y_box, bw, bh, color=GRAY)
    t(s, step, x, y_box, bw, bh, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        t(s, "→", x+bw+0.01, y_box+0.12, 0.25, 0.4, size=18, color=MUTED, align=PP_ALIGN.CENTER)

hline(s, 2.2, color=GRAY)

# Key facts
label(s, "重要特性", 0.55, 2.3)
facts = [
    "單次 LLM 推理可在同一個 content 陣列中同時產生 text + 多個 tool_use",
    "多個 tool_use 批次發出 → 平行執行 → 所有結果一起回來 → 才進行下一次推理",
    "Loop 終止條件：LLM 推理產生的回應不含任何 tool_use，控制權回到使用者",
    "Background Bash（run_in_background）不阻塞推理，結果以非同步通知送達",
    "Subagent 有自己獨立的 agentic loop，不佔主 agent 的推理次數",
]
bullets(s, facts, 0.55, 2.72, 12.2, size=17, gap=0.52)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4a – 工具總覽 (1/2)  Pre-loaded 工具
# ══════════════════════════════════════════════════════════════════
TEAL = RGBColor(0x0E, 0x7A, 0x8E)

s = prs.slides.add_slide(BLANK)
slide_header(s, "02  工具總覽 (1/2) — Pre-loaded 工具（9 個）",
             "Session 開始即載入 schema，可直接呼叫，零延遲")

# (color, name, perm, description)
preloaded = [
    (BLUE,   "Read",       "auto-allow", "讀取本地檔案；支援圖片（base64）、PDF（分頁）、ipynb（解析 cells）；硬上限 25k tokens"),
    (BLUE,   "Glob",       "auto-allow", "用 glob pattern 搜尋符合條件的檔案路徑，結果依修改時間排序"),
    (BLUE,   "Grep",       "auto-allow", "直接 spawn ripgrep 搜尋檔案內容（bypass shell，無注入風險，支援 regex）"),
    (MUTED,  "ToolSearch", "auto-allow", "將 deferred 工具的 schema 載入 context，讓 Claude 能夠呼叫它"),
    (GREEN,  "Edit",       "需確認",     "精確 old_string→new_string 替換；old_string 必須唯一；只傳 diff，token 消耗最少"),
    (GREEN,  "Write",      "需確認",     "覆寫整個檔案內容；適合新建檔案或完全重寫；自動計算 diff 顯示給使用者"),
    (ORANGE, "Bash",       "需確認",     "執行任意 shell 命令；支援 run_in_background（完成後以 task-notification 通知）"),
    (MUTED,  "Agent",      "需確認",     "啟動有獨立 context 的 subagent（可用不同 model）執行複雜子任務"),
    (MUTED,  "Skill",      "需確認",     "呼叫 .claude/skills/ 中定義的自訂指令集；LLM 接收展開後的內容再推理"),
]

# column x positions and widths
NX, NW = 0.55, 2.05   # name
PX, PW = 2.65, 1.35   # perm badge
DX, DW = 4.08, 9.1    # description
ROW_H  = 0.56

cy = 1.22
rect(s, 0.5, cy, 12.33, 0.40, color=GRAY)
t(s, "工具",       NX+0.05, cy+0.07, NW,  0.30, size=13, bold=True, color=MUTED)
t(s, "權限",       PX+0.05, cy+0.07, PW,  0.30, size=13, bold=True, color=MUTED)
t(s, "一句說明",   DX+0.05, cy+0.07, DW,  0.30, size=13, bold=True, color=MUTED)
cy += 0.42

for col, name, perm, desc in preloaded:
    hline(s, cy, color=GRAY)
    t(s, name,  NX+0.05, cy+0.07, NW,  ROW_H-0.1, size=15, bold=True,  color=col)
    pc = BLUE if perm == "auto-allow" else MUTED
    t(s, perm,  PX+0.05, cy+0.07, PW,  ROW_H-0.1, size=12, color=pc,   italic=True)
    t(s, desc,  DX+0.05, cy+0.07, DW,  ROW_H-0.1, size=14, color=INK)
    cy += ROW_H

hline(s, cy, color=GRAY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4b – 工具總覽 (2/2)  Deferred + MCP
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "02  工具總覽 (2/2) — Deferred（16 個）+ MCP（2 個）",
             "Deferred 需先 ToolSearch 載入 schema（+1 round-trip）；MCP 由外部 server 提供")

deferred = [
    # (color, name, group, description)
    (MUTED,  "AskUserQuestion", "互動",   "主動向使用者提問並等待回答（唯一能讓 Claude 暫停等待輸入的工具）"),
    (BLUE,   "WebFetch",        "唯讀",   "抓取指定 URL 的網頁內容，HTML 自動轉 markdown 後回傳"),
    (BLUE,   "WebSearch",       "唯讀",   "搜尋網路，取得超出 training data 的最新資訊"),
    (GREEN,  "NotebookEdit",    "寫入",   "編輯 Jupyter notebook cell（.ipynb 必須用此；直接 Edit 會報錯）"),
    (PURPLE, "EnterPlanMode",   "規劃",   "進入 Plan 模式；計畫存至 ~/.claude/plans/{slug}.md"),
    (PURPLE, "ExitPlanMode",    "規劃",   "離開 Plan 模式，提交計畫給使用者審核"),
    (PURPLE, "EnterWorktree",   "隔離",   "在 .claude/worktrees/ 建立隔離的 git worktree 環境"),
    (RED,    "TaskCreate",      "任務",   "建立新的使用者任務（ID 為數字，如 1）"),
    (RED,    "TaskUpdate",      "任務",   "更新任務的狀態或內容"),
    (RED,    "TaskStop",        "任務",   "停止指定的背景 Bash 任務（ID 為 hash，如 bz6r...；無法停 TaskCreate 任務）"),
    (RED,    "CronCreate",      "排程",   "建立定時排程（session-only，Claude 退出後消失）"),
    (RED,    "CronDelete",      "排程",   "刪除指定排程"),
    (BLUE,   "TaskGet",         "唯讀",   "取得指定任務的詳細資訊"),
    (BLUE,   "TaskList",        "唯讀",   "列出所有使用者任務"),
    (BLUE,   "TaskOutput",      "唯讀",   "讀取背景 Bash 任務的標準輸出"),
    (BLUE,   "CronList",        "唯讀",   "列出目前 session 中所有排程任務"),
]
mcp_tools = [
    (TEAL,   "mcp__ide__getDiagnostics", "MCP", "取得 IDE（VS Code / JetBrains）的 lint／type 錯誤診斷資訊"),
    (TEAL,   "mcp__ide__executeCode",    "MCP", "在 IDE 環境執行指定程式碼片段"),
]

NX2, NW2 = 0.50, 2.60
GX2, GW2 = 3.15, 1.00
DX2, DW2 = 4.22, 8.95
ROW_H2 = 0.39

cy = 1.22
rect(s, 0.5, cy, 12.33, 0.36, color=GRAY)
t(s, "工具",     NX2+0.05, cy+0.06, NW2, 0.28, size=12, bold=True, color=MUTED)
t(s, "群組",     GX2+0.05, cy+0.06, GW2, 0.28, size=12, bold=True, color=MUTED)
t(s, "一句說明", DX2+0.05, cy+0.06, DW2, 0.28, size=12, bold=True, color=MUTED)
cy += 0.38

for col, name, grp, desc in deferred:
    hline(s, cy, color=GRAY)
    t(s, name, NX2+0.05, cy+0.05, NW2, ROW_H2, size=13, bold=True, color=col)
    t(s, grp,  GX2+0.05, cy+0.05, GW2, ROW_H2, size=11, color=MUTED, italic=True)
    t(s, desc, DX2+0.05, cy+0.05, DW2, ROW_H2, size=13, color=INK)
    cy += ROW_H2

# MCP divider
hline(s, cy, color=BLUE, h=0.06)
t(s, "MCP（需外部 server 連線）", NX2+0.05, cy+0.08, 4, 0.30, size=12, bold=True, color=TEAL)
cy += 0.38
for col, name, grp, desc in mcp_tools:
    hline(s, cy, color=GRAY)
    t(s, name, NX2+0.05, cy+0.05, NW2, ROW_H2, size=13, bold=True, color=col)
    t(s, grp,  GX2+0.05, cy+0.05, GW2, ROW_H2, size=11, color=TEAL, italic=True)
    t(s, desc, DX2+0.05, cy+0.05, DW2, ROW_H2, size=13, color=INK)
    cy += ROW_H2
hline(s, cy, color=GRAY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 – Claude Code 怎麼寫 Code
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "03  Claude Code 怎麼寫 Code — 流程與設計優勢",
             "為何 Claude Code 是當前最強 code agent")

# Flow + explanation
steps_code = [
    ("①  探索",  BLUE,   "Glob / Grep / Read",
     [
         "Glob：pattern 搜尋，直接拿到目標路徑，不走 shell",
         "Grep：直接 spawn ripgrep binary，bypass shell，不需 Bash 確認",
         "Read：預載、預設 auto-allow，不需使用者確認，直接讀，速度最快",
         "三者都是唯讀，預設自動允許 → 零中斷探索整個 codebase",
     ]),
    ("②  修改",  GREEN,  "Edit / Write",
     [
         "Edit：只傳 old_string→new_string diff，token 消耗極少",
         "old_string 唯一性強制驗證 → 確保修改位置精準，不會誤改",
         "Write：整個檔案覆寫，適合新建或完全重寫，自動產生 diff 給使用者確認",
         "兩者都鎖定工作目錄範圍，不會意外修改到其他地方",
     ]),
    ("③  驗證",  ORANGE, "Bash",
     [
         "跑測試、build、lint，確認修改正確",
         "run_in_background：長時間命令不阻塞，Claude 可繼續其他推理",
         "structured exit code：0 成功，非 0 觸發 Claude 自動分析錯誤並修正",
     ]),
]

cy = 1.28
for step_label, col, tools_str, detail_bullets in steps_code:
    rect(s, 0.5, cy, 12.33, 0.36, color=GRAY)
    t(s, f"{step_label}  ·  {tools_str}", 0.62, cy+0.03, 12, 0.32,
      size=15, bold=True, color=col)
    for item in detail_bullets:
        t(s, f"    • {item}", 0.62, cy+0.38, 12, 0.38, size=16, color=INK)
        cy += 0.38
    cy += 0.52

hline(s, cy+0.1, color=GRAY)
t(s, "結果：每一步工具都針對「最少確認次數 + 最少 token 消耗 + 最精準操作」最佳化",
  0.55, cy+0.18, 12.2, 0.42, size=16, bold=True, color=BLUE)


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 – Pre-loaded vs Deferred + ToolSearch
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  Tool 特性 — Pre-loaded vs Deferred vs MCP")

# Table-like layout
headers = ["分類", "數量", "Schema 載入時機", "代表工具", "特性"]
col_x   = [0.5,   1.8,   3.0,               5.5,        9.5]
col_w   = [1.2,   1.1,   2.4,               3.8,        3.5]

# header row
rect(s, 0.5, 1.3, 12.33, 0.45, color=GRAY)
for hdr, x, w in zip(headers, col_x, col_w):
    t(s, hdr, x+0.06, 1.33, w, 0.38, size=15, bold=True, color=MUTED)

rows = [
    (BLUE,   "Pre-loaded",  "9",  "Session 開始即載入",
     "Read, Edit, Glob, Grep, Bash, Agent, Skill, Write, ToolSearch",
     "零延遲，直接呼叫"),
    (GREEN,  "Deferred",    "16", "需先呼叫 ToolSearch",
     "WebFetch, WebSearch, AskUserQuestion, NotebookEdit, EnterPlanMode...",
     "省 token，多一次 round-trip（~2-3 秒）"),
    (ORANGE, "MCP",         "不定","Server 連線時載入",
     "mcp__ide__getDiagnostics, mcp__ide__executeCode, ...",
     "由外部 server 提供，未連線時不存在"),
]

cy = 1.8
for col, cat, count, timing, tools_str, note in rows:
    hline(s, cy, color=GRAY)
    vals = [cat, count, timing, tools_str, note]
    colors = [col, INK, INK, INK, MUTED]
    for v, c, x, w in zip(vals, colors, col_x, col_w):
        t(s, v, x+0.06, cy+0.1, w, 0.5, size=14, color=c, bold=(c==col))
    cy += 0.62
hline(s, cy, color=GRAY)

label(s, "ToolSearch 的角色", 0.5, cy+0.15)
bullets(s, [
    "唯一預載的「載入其他工具」的 meta 工具",
    "用法：ToolSearch(\"select:WebFetch\") 或 ToolSearch(\"list directory\")（關鍵字搜尋）",
    "同一 session 內載入一次即可，不需重複呼叫；跨 session 需重新載入",
], 0.5, cy+0.52, 12.2, size=16, gap=0.48)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 – Read vs Bash / Edit vs Write
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  Tool 特性 — Read / Edit / Write vs Bash 的本質差異")

label(s, "為什麼有 Read/Edit/Write，不直接用 Bash？", 0.5, 1.15, color=BLUE)

headers2 = ["面向", "Read / Write / Edit", "Bash"]
col_x2 = [0.5, 3.7, 9.0]
col_w2 = [3.1, 5.2, 4.1]
rect(s, 0.5, 1.55, 12.33, 0.42, color=GRAY)
for hdr, x, w in zip(headers2, col_x2, col_w2):
    t(s, hdr, x+0.08, 1.58, w, 0.36, size=15, bold=True, color=MUTED)

rows2 = [
    ("預設權限",   "Read auto-allow；Write/Edit 分開設定",         "全部需要 prefix / exact / wildcard 規則"),
    ("匹配邏輯",   "filePatternTools（路徑型）",                    "bashPrefixTools（命令型）"),
    ("Token 限制", "Read：硬上限 25,000 tokens，超過回 error",       "無強制限制"),
    ("Diff 顯示",  "自動計算 structuredPatch + gitDiff 給使用者看",  "無"),
    ("寫入範圍",   "鎖定工作目錄及子目錄（框架層強制）",              "不限制"),
    ("特殊格式",   "圖片（base64）/ .ipynb（解析 cells）/ PDF（分頁）","只有純文字 / bytes"),
    ("注入風險",   "無，參數由 framework 組裝",                      "命令字串經過 shell 解析，有注入風險"),
]

cy = 2.0
for dim, left, right in rows2:
    hline(s, cy, color=GRAY)
    t(s, dim,   col_x2[0]+0.08, cy+0.08, col_w2[0]-0.1, 0.42, size=15, color=MUTED, bold=True)
    t(s, left,  col_x2[1]+0.08, cy+0.08, col_w2[1]-0.1, 0.42, size=14, color=BLUE)
    t(s, right, col_x2[2]+0.08, cy+0.08, col_w2[2]-0.1, 0.42, size=14, color=ORANGE)
    cy += 0.56
hline(s, cy, color=GRAY)

label(s, "Edit vs Write", 0.5, cy+0.12, color=GREEN)
t(s, "Edit：精確字串替換（old→new），只傳 diff，token 少，old_string 必須唯一否則報錯，日常修改首選",
  0.5, cy+0.5, 12.2, 0.42, size=16, color=INK)
t(s, "Write：整個檔案覆寫，適合新建檔案或完全重寫，token 消耗多",
  0.5, cy+0.95, 12.2, 0.4, size=16, color=INK)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 – Skill System
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  Tool 特性 — Skill 系統")

label(s, "Skill 是什麼", 0.5, 1.15, color=BLUE)
bullets(s, [
    "放在 .claude/skills/<name>/SKILL.md 的 Markdown 指令檔，告訴 Claude 如何完成特定任務",
    "可附帶腳本，在展開給 LLM 之前先預處理資料（! 前綴命令）",
], 0.5, 1.52, 12.2, size=17, gap=0.5)

hline(s, 2.56, color=GRAY)
label(s, "兩種觸發方式", 0.5, 2.65, color=BLUE)

rect(s, 0.5, 3.05, 5.9, 0.38, color=GRAY)
t(s, "Trigger A  /skill-name  （使用者輸入）", 0.62, 3.08, 5.65, 0.32, size=15, bold=True, color=BLUE)
bullets(s, [
    "CLI 在框架層讀取 SKILL.md，展開為 isMeta:true 訊息注入 context",
    "Claude 不需呼叫 Skill tool，直接看到展開後的指令",
], 0.62, 3.46, 5.65, size=16, gap=0.46)

rect(s, 6.9, 3.05, 5.93, 0.38, color=GRAY)
t(s, "Trigger B  Skill tool  （Claude 主動呼叫）", 7.02, 3.08, 5.7, 0.32, size=15, bold=True, color=PURPLE)
bullets(s, [
    "Claude 先 ToolSearch 載入 Skill schema，再呼叫 Skill(\"name\", args)",
    "Framework 處理 SKILL.md → tool_result 回傳給 LLM",
], 7.02, 3.46, 5.7, size=16, gap=0.46)

hline(s, 4.44, color=GRAY)
label(s, "重要設計差異：allowed-tools vs Agent 工具隔離", 0.5, 4.52, color=RED)
bullets(s, [
    "Skill allowed-tools = 軟性：自動允許列出的工具，但 LLM 仍可呼叫其他工具（需使用者確認）",
    "Agent tools = 硬性：API 層隔離，LLM 根本看不到白名單以外的工具 schema",
    "! 前綴命令：Framework 執行後再注入，LLM 看不到原始命令，用於大量資料預處理",
], 0.5, 4.92, 12.2, size=16, gap=0.5)


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 – Agent Tool
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  Tool 特性 — Agent Tool（Subagent）")

label(s, "5 種 subagent_type", 0.5, 1.15, color=BLUE)

types = [
    ("general-purpose", "全部工具（*）",                               "通用型，複雜多步驟任務"),
    ("Explore",         "除 Agent, ExitPlanMode, Edit, Write 等",      "快速探索 codebase；3 種深度（quick/medium/very thorough）"),
    ("Plan",            "同 Explore",                                  "架構規劃，回傳 step-by-step 實作計畫"),
    ("claude-code-guide","Glob, Grep, Read, WebFetch, WebSearch",      "回答 Claude Code / SDK / API 問題；支援 resume 繼續"),
    ("statusline-setup","Read, Edit",                                  "設定 Claude Code status line"),
]
headers3 = ["subagent_type", "可用工具", "用途"]
col_x3 = [0.5, 3.4, 7.4]
col_w3 = [2.8, 3.9, 5.2]
rect(s, 0.5, 1.52, 12.33, 0.4, color=GRAY)
for h, x, w in zip(headers3, col_x3, col_w3):
    t(s, h, x+0.08, 1.55, w, 0.34, size=14, bold=True, color=MUTED)
cy = 1.95
for name, tools_str, desc in types:
    hline(s, cy, color=GRAY)
    t(s, name, col_x3[0]+0.08, cy+0.08, col_w3[0], 0.42, size=14, bold=True, color=BLUE)
    t(s, tools_str, col_x3[1]+0.08, cy+0.08, col_w3[1], 0.42, size=13, color=INK)
    t(s, desc, col_x3[2]+0.08, cy+0.08, col_w3[2], 0.42, size=13, color=INK)
    cy += 0.55
hline(s, cy, color=GRAY)

label(s, "關鍵特性", 0.5, cy+0.12, color=BLUE)
bullets(s, [
    "獨立 context：每個 subagent 有自己的 context window，不共享主 agent 對話歷史",
    "不同 model：實測 Explore 使用 claude-haiku（快/省錢），主 agent 使用 claude-sonnet",
    "結果包含 totalDurationMs / totalTokens / totalToolUseCount — 完整成本追蹤",
    "agentId 可 resume：下次呼叫帶入 resume 參數可繼續同一個 subagent",
], 0.5, cy+0.5, 12.2, size=16, gap=0.48)


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 – Hooks
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "04  Tool 特性 — Hooks（事件驅動擴充）")

label(s, "21 種支援事件（來源：binary 2.1.72）", 0.5, 1.15, color=BLUE)

events = [
    "PreToolUse / PostToolUse / PostToolUseFailure",
    "UserPromptSubmit / Notification",
    "SessionStart / SessionEnd / Stop",
    "SubagentStart / SubagentStop / TeammateIdle",
    "PreCompact / PermissionRequest / Setup",
    "TaskCompleted / Elicitation / ElicitationResult",
    "ConfigChange / WorktreeCreate / WorktreeRemove / InstructionsLoaded",
]
cy = 1.52
for line in events:
    t(s, "• " + line, 0.6, cy, 12.1, 0.4, size=15, color=INK)
    cy += 0.4

hline(s, cy+0.08, color=GRAY)
label(s, "兩種 Hook 類型", 0.5, cy+0.2, color=BLUE)

# command type
rect(s, 0.5, cy+0.58, 5.9, 0.38, color=GRAY)
t(s, "command type  （執行 shell 命令）", 0.62, cy+0.61, 5.65, 0.32, size=15, bold=True, color=GREEN)
cmd = [
    "command（必填）：執行的 shell 命令",
    "async：背景執行，不阻塞",
    "asyncRewake：背景執行，exit code 2 時喚醒 model",
    "once：執行一次後自動移除",
    "timeout / statusMessage",
]
cy2 = cy+1.0
for item in cmd:
    t(s, "  • " + item, 0.62, cy2, 5.65, 0.38, size=15, color=INK)
    cy2 += 0.38

# prompt type
rect(s, 6.9, cy+0.58, 5.93, 0.38, color=GRAY)
t(s, "prompt type  （LLM 評估）", 7.02, cy+0.61, 5.7, 0.32, size=15, bold=True, color=PURPLE)
prm = [
    "prompt（必填）：提示文字，$ARGUMENTS 取 hook 輸入",
    "model：指定 model（預設用小型快速 model）",
    "once / timeout / statusMessage",
]
cy3 = cy+1.0
for item in prm:
    t(s, "  • " + item, 7.02, cy3, 5.7, 0.38, size=15, color=INK)
    cy3 += 0.38


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 – Permission System
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "05  Permission 系統")

label(s, "兩個層次", 0.5, 1.15, color=BLUE)

rect(s, 0.5, 1.52, 12.33, 0.38, color=GRAY)
t(s, "Layer 1  工具權限（可設定）", 0.62, 1.55, 6, 0.32, size=15, bold=True, color=BLUE)
bullets(s, [
    'permissions.allow / deny 設定於 settings.json',
    '格式：工具名稱 "Edit"，帶參數 "Bash(git:*)"，domain "WebFetch(domain:docs.anthropic.com)"',
], 0.62, 1.93, 12.1, size=16, gap=0.46)

rect(s, 0.5, 2.9, 12.33, 0.38, color=GRAY)
t(s, "Layer 2  判斷性確認（不可關閉）", 0.62, 2.93, 7, 0.32, size=15, bold=True, color=RED)
bullets(s, [
    "Claude 自行判斷：操作超出請求範圍、發現非預期狀態、請求模糊時會主動暫停",
    "無法透過設定關閉，這是 Claude 本身的風險評估，不是規則系統",
], 0.62, 3.31, 12.1, size=16, gap=0.46)

hline(s, 4.28, color=GRAY)
label(s, "Bash 匹配邏輯（來源：binary 函數逆向）", 0.5, 4.36, color=BLUE)
bash_items = [
    'prefix  Bash(git:*)  →  匹配 "git"、"git status"、"xargs git"',
    'exact   Bash(git status)  →  只匹配完全相符的字串',
    'wildcard  Bash(git *)  →  glob 模式匹配',
    '安全設計：複合命令（&&, ||, ;, |）永遠不匹配 prefix 規則 → 防止 "git status && rm -rf /" 越獄',
]
bullets(s, bash_items, 0.5, 4.74, 12.2, size=16, gap=0.5)

hline(s, 6.82, color=GRAY)
label(s, "Permission Modes", 0.5, 6.88, color=BLUE, size=13)
modes = [("default","預設，需確認才執行"),("acceptEdits","自動接受檔案編輯，Bash 仍需確認"),
         ("bypassPermissions","跳過所有確認（危險）"),("dontAsk","完全不詢問"),
         ("plan","僅規劃，限制實際執行")]
cx = 0.55
for mode, note in modes:
    t(s, mode, cx, 7.18, 2.1, 0.3, size=12, bold=True, color=BLUE)
    t(s, note, cx+2.15, 7.18, 2.2, 0.3, size=12, color=MUTED)
    cx += 4.35
    if cx > 9:
        break


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 – Takeaways
# ══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_header(s, "06  Takeaways")

takes = [
    (BLUE,   "Agentic Loop",
     "批次 tool call 平行執行，不含 tool_use 才停止。Subagent 有獨立 loop，不佔主 agent 的推理次數。"),
    (GREEN,  "工具設計哲學",
     "讀取工具預載且 auto-allow → 零中斷探索。Edit 只傳 diff → token 最少。"
     "每個工具都針對最少確認次數、最精準操作最佳化。"),
    (ORANGE, "Pre-loaded vs Deferred",
     "9 個預載（零延遲）vs 16 個延遲載入（省 token，透過 ToolSearch，+2-3 秒）。"
     "根據使用頻率權衡 latency 與 token budget。"),
    (PURPLE, "Skill & Agent 的隔離差異",
     "Skill allowed-tools 是軟性（context 共享，只是自動允許）。"
     "Agent tools 是硬性 API 層隔離（LLM 根本看不到其他工具 schema）。"),
    (RED,    "Permission 安全設計",
     "複合命令永遠不匹配 prefix 規則（binary 逆向確認），防止越獄。"
     "Layer 2 判斷性確認無法用設定關閉。"),
    (MUTED,  "觀察方式",
     "所有機制都可透過三種方式交叉驗證："
     "真實 Session JSONL（runtime ground truth）、"
     "Claude Code binary 原始碼（函數邏輯）、官網文件（設計意圖）。"),
]

cy = 1.28
for col, title, body in takes:
    hline(s, cy, color=GRAY)
    t(s, title, 0.55, cy+0.1, 2.6, 0.48, size=17, bold=True, color=col)
    t(s, body,  3.25, cy+0.1, 9.55, 0.55, size=16, color=INK)
    cy += 0.82
hline(s, cy, color=GRAY)


# ── Save ──────────────────────────────────────────────────────
out = "/home/user/reseach_claude_code/agent_architecture.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
